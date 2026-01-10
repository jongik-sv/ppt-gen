#!/usr/bin/env python3
"""
Create thumbnail grids from PowerPoint presentation slides.

Creates a grid layout of slide thumbnails with configurable columns (max 6).
Each grid contains up to cols×(cols+1) images. For presentations with more
slides, multiple numbered grid files are created automatically.

The program outputs the names of all files created.

Output:
- Single grid: {prefix}.jpg (if slides fit in one grid)
- Multiple grids: {prefix}-1.jpg, {prefix}-2.jpg, etc.

Grid limits by column count:
- 3 cols: max 12 slides per grid (3×4)
- 4 cols: max 20 slides per grid (4×5)
- 5 cols: max 30 slides per grid (5×6) [default]
- 6 cols: max 42 slides per grid (6×7)

Usage:
    python thumbnail.py input.pptx [output_prefix] [--cols N] [--outline-placeholders]
    python thumbnail.py input.pptx output/ --slides 5 --single
    python thumbnail.py --from-images images/ output/ [--cols N]

Examples:
    python thumbnail.py presentation.pptx
    # Creates: thumbnails.jpg (using default prefix)
    # Outputs:
    #   Created 1 grid(s):
    #     - thumbnails.jpg

    python thumbnail.py large-deck.pptx grid --cols 4
    # Creates: grid-1.jpg, grid-2.jpg, grid-3.jpg
    # Outputs:
    #   Created 3 grid(s):
    #     - grid-1.jpg
    #     - grid-2.jpg
    #     - grid-3.jpg

    python thumbnail.py template.pptx analysis --outline-placeholders
    # Creates thumbnail grids with red outlines around text placeholders

    python thumbnail.py input.pptx output/ --slides 5 --single
    # Creates: output/slide-5.png (single slide at 1980x1080)
    # Outputs:
    #   Created single thumbnail:
    #     - output/slide-5.png

    python thumbnail.py input.pptx output/ --slides 1,3,5
    # Creates: output/slide-1.png, output/slide-3.png, output/slide-5.png

    # Image-based mode (no LibreOffice required):
    python thumbnail.py --from-images images/ grid
    # Creates grid from existing images in the directory

    python thumbnail.py --from-images images/ output/ --single
    # Resizes images to 1980x1080 thumbnails
"""

# Single slide output constants
SINGLE_WIDTH = 1980
SINGLE_HEIGHT = 1080

import argparse
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from inventory import extract_text_inventory
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation


# Office backend detection cache
_office_backend_cache = None


def find_office_backend() -> str:
    """사용 가능한 Office 백엔드 찾기 (LibreOffice 우선)

    Returns:
        'libreoffice' | 'msoffice' | None
    """
    global _office_backend_cache
    if _office_backend_cache is not None:
        return _office_backend_cache

    # 1. LibreOffice 확인 (우선)
    if shutil.which('soffice'):
        _office_backend_cache = 'libreoffice'
        return 'libreoffice'

    # 2. MS Office 확인 (Windows 전용)
    if sys.platform == 'win32':
        try:
            import win32com.client
            ppt = win32com.client.Dispatch('PowerPoint.Application')
            ppt.Quit()
            _office_backend_cache = 'msoffice'
            return 'msoffice'
        except Exception:
            pass

    _office_backend_cache = None
    return None


def convert_with_libreoffice(pptx_path: Path, output_dir: Path) -> Path:
    """LibreOffice로 PPTX→PDF 변환"""
    result = subprocess.run(
        [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(pptx_path),
        ],
        capture_output=True,
        text=True,
    )

    pdf_path = output_dir / f"{pptx_path.stem}.pdf"
    if result.returncode != 0 or not pdf_path.exists():
        raise RuntimeError(f"LibreOffice PDF 변환 실패: {result.stderr}")

    return pdf_path


def convert_with_msoffice(pptx_path: Path, output_dir: Path) -> Path:
    """MS PowerPoint COM으로 PPTX→PDF 변환 (Windows 전용)"""
    import win32com.client
    import pythoncom

    # COM 초기화
    pythoncom.CoInitialize()

    ppt = None
    presentation = None
    pdf_path = output_dir / f"{pptx_path.stem}.pdf"

    try:
        ppt = win32com.client.Dispatch('PowerPoint.Application')
        ppt.Visible = True  # 백그라운드 실행 제약으로 인해 필요

        # 절대 경로를 Windows 스타일로 변환
        pptx_abs_path = str(pptx_path.absolute()).replace('/', '\\')
        pdf_abs_path = str(pdf_path.absolute()).replace('/', '\\')

        presentation = ppt.Presentations.Open(
            pptx_abs_path,
            ReadOnly=True,
            WithWindow=False
        )

        # ppFixedFormatTypePDF = 2, ppFixedFormatIntentPrint = 2
        presentation.SaveAs(pdf_abs_path, 32)  # ppSaveAsPDF = 32

        presentation.Close()
        presentation = None

    finally:
        if presentation:
            try:
                presentation.Close()
            except:
                pass
        if ppt:
            try:
                ppt.Quit()
            except:
                pass
        pythoncom.CoUninitialize()

    if not pdf_path.exists():
        raise RuntimeError("MS Office PDF 변환 실패")

    return pdf_path


def convert_pptx_to_pdf(pptx_path: Path, output_dir: Path, backend: str = None) -> Path:
    """PPTX→PDF 변환 (자동 백엔드 선택)

    Args:
        pptx_path: 입력 PPTX 파일
        output_dir: 출력 디렉토리
        backend: 'libreoffice' | 'msoffice' | 'auto' | None

    Returns:
        생성된 PDF 파일 경로
    """
    if backend is None or backend == 'auto':
        backend = find_office_backend()

    if backend is None:
        raise RuntimeError(
            "LibreOffice 또는 MS Office가 필요합니다.\n"
            "- LibreOffice: https://www.libreoffice.org/\n"
            "- MS Office: PowerPoint가 설치되어 있어야 합니다 (Windows)"
        )

    if backend == 'libreoffice':
        return convert_with_libreoffice(pptx_path, output_dir)
    else:
        return convert_with_msoffice(pptx_path, output_dir)

# Constants
THUMBNAIL_WIDTH = 300  # Fixed thumbnail width in pixels
CONVERSION_DPI = 100  # DPI for PDF to image conversion
MAX_COLS = 6  # Maximum number of columns
DEFAULT_COLS = 5  # Default number of columns
JPEG_QUALITY = 95  # JPEG compression quality

# Grid layout constants
GRID_PADDING = 20  # Padding between thumbnails
BORDER_WIDTH = 2  # Border width around thumbnails
FONT_SIZE_RATIO = 0.12  # Font size as fraction of thumbnail width
LABEL_PADDING_RATIO = 0.4  # Label padding as fraction of font size


def main():
    parser = argparse.ArgumentParser(
        description="Create thumbnail grids from PowerPoint slides."
    )
    parser.add_argument("input", nargs="?", help="Input PowerPoint file (.pptx)")
    parser.add_argument(
        "output_prefix",
        nargs="?",
        default="thumbnails",
        help="Output prefix for image files (default: thumbnails, will create prefix.jpg or prefix-N.jpg)",
    )
    parser.add_argument(
        "--cols",
        type=int,
        default=DEFAULT_COLS,
        help=f"Number of columns (default: {DEFAULT_COLS}, max: {MAX_COLS})",
    )
    parser.add_argument(
        "--outline-placeholders",
        action="store_true",
        help="Outline text placeholders with a colored border",
    )
    parser.add_argument(
        "--slides",
        type=str,
        help="Specific slide(s) to extract (0-based index). Use comma for multiple: --slides 5 or --slides 1,3,5",
    )
    parser.add_argument(
        "--single",
        action="store_true",
        help="Output single slide image(s) instead of grid (1980x1080 PNG)",
    )
    parser.add_argument(
        "--from-images",
        type=str,
        metavar="DIR",
        help="Create thumbnails from existing images in directory (no LibreOffice needed)",
    )
    parser.add_argument(
        "--prefix",
        type=str,
        default="slide",
        help="Prefix for output filenames (default: slide). Output: {prefix}-{index}.png",
    )
    parser.add_argument(
        "--backend",
        type=str,
        choices=["auto", "libreoffice", "msoffice"],
        default="auto",
        help="PDF conversion backend (default: auto). auto=LibreOffice first, then MS Office",
    )

    args = parser.parse_args()

    # Show detected backend
    if not args.from_images:
        backend = args.backend if args.backend != 'auto' else find_office_backend()
        if backend:
            print(f"Using backend: {backend}")
        else:
            print("Error: No Office application found (LibreOffice or MS Office required)")
            sys.exit(1)

    # Handle image-based mode
    if args.from_images:
        handle_from_images_mode(args)
        return

    # Validate columns
    cols = min(args.cols, MAX_COLS)
    if args.cols > MAX_COLS:
        print(f"Warning: Columns limited to {MAX_COLS} (requested {args.cols})")

    # Validate input
    input_path = Path(args.input)
    if not input_path.exists() or input_path.suffix.lower() != ".pptx":
        print(f"Error: Invalid PowerPoint file: {args.input}")
        sys.exit(1)

    # Parse slide indices if specified
    slide_indices = None
    if args.slides:
        try:
            slide_indices = [int(s.strip()) for s in args.slides.split(",")]
        except ValueError:
            print(f"Error: Invalid slide indices: {args.slides}")
            sys.exit(1)

    # Handle single slide mode
    if args.single or slide_indices:
        output_dir = Path(args.output_prefix)
        output_dir.mkdir(parents=True, exist_ok=True)
        create_single_thumbnails(input_path, output_dir, slide_indices, args.single, args.backend)
        return

    # Construct output path (always JPG for grid mode)
    output_path = Path(f"{args.output_prefix}.jpg")

    print(f"Processing: {args.input}")

    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            # Get placeholder regions if outlining is enabled
            placeholder_regions = None
            slide_dimensions = None
            if args.outline_placeholders:
                print("Extracting placeholder regions...")
                placeholder_regions, slide_dimensions = get_placeholder_regions(
                    input_path
                )
                if placeholder_regions:
                    print(f"Found placeholders on {len(placeholder_regions)} slides")

            # Convert slides to images
            slide_images = convert_to_images(input_path, Path(temp_dir), CONVERSION_DPI, args.backend)
            if not slide_images:
                print("Error: No slides found")
                sys.exit(1)

            print(f"Found {len(slide_images)} slides")

            # Create grids (max cols×(cols+1) images per grid)
            grid_files = create_grids(
                slide_images,
                cols,
                THUMBNAIL_WIDTH,
                output_path,
                placeholder_regions,
                slide_dimensions,
            )

            # Print saved files
            print(f"Created {len(grid_files)} grid(s):")
            for grid_file in grid_files:
                print(f"  - {grid_file}")

    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)


def create_hidden_slide_placeholder(size):
    """Create placeholder image for hidden slides."""
    img = Image.new("RGB", size, color="#F0F0F0")
    draw = ImageDraw.Draw(img)
    line_width = max(5, min(size) // 100)
    draw.line([(0, 0), size], fill="#CCCCCC", width=line_width)
    draw.line([(size[0], 0), (0, size[1])], fill="#CCCCCC", width=line_width)
    return img


def handle_from_images_mode(args):
    """Handle thumbnail generation from existing images (no LibreOffice needed).

    This mode:
    1. Reads images from a directory
    2. Resizes them to create thumbnails
    3. Creates a grid or individual thumbnails
    """
    image_dir = Path(args.from_images)
    if not image_dir.exists():
        print(f"Error: Image directory not found: {image_dir}")
        sys.exit(1)

    # Find all images in directory (sorted by name)
    image_extensions = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"}
    image_files = sorted([
        f for f in image_dir.iterdir()
        if f.suffix.lower() in image_extensions
    ])

    if not image_files:
        print(f"Error: No images found in {image_dir}")
        sys.exit(1)

    print(f"Found {len(image_files)} images in {image_dir}")

    # Validate columns
    cols = min(args.cols, MAX_COLS)
    if args.cols > MAX_COLS:
        print(f"Warning: Columns limited to {MAX_COLS} (requested {args.cols})")

    # Parse slide indices if specified
    slide_indices = None
    if args.slides:
        try:
            slide_indices = [int(s.strip()) for s in args.slides.split(",")]
        except ValueError:
            print(f"Error: Invalid slide indices: {args.slides}")
            sys.exit(1)

    # Handle single mode - create individual resized thumbnails
    if args.single or slide_indices:
        output_dir = Path(args.output_prefix) if args.output_prefix != "thumbnails" else Path(".")
        output_dir.mkdir(parents=True, exist_ok=True)
        create_thumbnails_from_images(image_files, output_dir, slide_indices, args.single, args.prefix)
        return

    # Grid mode - create thumbnail grid
    output_path = Path(f"{args.output_prefix}.jpg")
    print(f"Creating grid with {cols} columns...")

    # Create grids
    grid_files = create_grids(
        image_files,
        cols,
        THUMBNAIL_WIDTH,
        output_path,
        placeholder_regions=None,
        slide_dimensions=None,
    )

    print(f"Created {len(grid_files)} grid(s):")
    for grid_file in grid_files:
        print(f"  - {grid_file}")


def create_thumbnails_from_images(image_files, output_dir, slide_indices=None, single_mode=False, prefix="slide"):
    """Create single thumbnails from existing images by resizing.

    Args:
        image_files: List of image file paths
        output_dir: Output directory
        slide_indices: List of indices to process (0-based). If None, process all.
        single_mode: If True, always create individual files
        prefix: Prefix for output filenames (default: "slide")
    """
    # Default to all images if not specified
    if slide_indices is None:
        slide_indices = list(range(len(image_files)))

    # Validate indices
    valid_indices = [i for i in slide_indices if 0 <= i < len(image_files)]
    if not valid_indices:
        print(f"Error: No valid indices (total images: {len(image_files)})")
        sys.exit(1)

    created_files = []
    for idx in valid_indices:
        input_path = image_files[idx]
        output_path = output_dir / f"{prefix}-{idx}.png"

        try:
            with Image.open(input_path) as img:
                # Convert to RGB if needed (for PNG with transparency)
                if img.mode in ("RGBA", "P"):
                    img = img.convert("RGB")

                # Resize to target dimensions maintaining aspect ratio with padding
                img_resized = resize_with_letterbox(img, SINGLE_WIDTH, SINGLE_HEIGHT)
                img_resized.save(output_path, "PNG")
                created_files.append(str(output_path))
                print(f"  Created: {output_path}")
        except Exception as e:
            print(f"  Warning: Failed to process {input_path}: {e}")

    # Print summary
    if len(created_files) == 1:
        print(f"Created single thumbnail:")
    else:
        print(f"Created {len(created_files)} thumbnail(s):")
    for f in created_files:
        print(f"  - {f}")


def resize_with_letterbox(img, target_width, target_height, bg_color=(255, 255, 255)):
    """Resize image to target dimensions with letterbox (padding) to maintain aspect ratio.

    Args:
        img: PIL Image
        target_width: Target width
        target_height: Target height
        bg_color: Background color for letterbox

    Returns:
        Resized PIL Image
    """
    orig_width, orig_height = img.size
    orig_aspect = orig_width / orig_height
    target_aspect = target_width / target_height

    if orig_aspect > target_aspect:
        # Image is wider - fit to width, pad top/bottom
        new_width = target_width
        new_height = int(target_width / orig_aspect)
    else:
        # Image is taller - fit to height, pad left/right
        new_height = target_height
        new_width = int(target_height * orig_aspect)

    # Resize
    img_resized = img.resize((new_width, new_height), Image.Resampling.LANCZOS)

    # Create canvas and paste centered
    canvas = Image.new("RGB", (target_width, target_height), bg_color)
    paste_x = (target_width - new_width) // 2
    paste_y = (target_height - new_height) // 2
    canvas.paste(img_resized, (paste_x, paste_y))

    return canvas


def create_single_thumbnails(pptx_path, output_dir, slide_indices=None, single_mode=False, backend=None):
    """Create single slide thumbnail(s) at high resolution.

    Args:
        pptx_path: Path to the PowerPoint file
        output_dir: Output directory for thumbnails
        slide_indices: List of slide indices (0-based) to extract. If None, extracts all.
        single_mode: If True and slide_indices has one item, outputs single file
        backend: 'libreoffice' | 'msoffice' | 'auto' | None
    """
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_dir = Path(temp_dir)

        # Get total slide count and hidden slides
        prs = Presentation(str(pptx_path))
        total_slides = len(prs.slides)
        hidden_slides = {
            idx for idx, slide in enumerate(prs.slides)
            if slide.element.get("show") == "0"
        }

        # Default to all slides if not specified
        if slide_indices is None:
            slide_indices = list(range(total_slides))

        # Validate indices
        valid_indices = [i for i in slide_indices if 0 <= i < total_slides]
        if not valid_indices:
            print(f"Error: No valid slide indices (total slides: {total_slides})")
            sys.exit(1)

        # Convert to PDF using unified backend
        print(f"Converting to PDF...")
        pdf_path = convert_pptx_to_pdf(pptx_path, temp_dir, backend)

        # Calculate DPI for target resolution (1980x1080)
        # Standard slide is 10" x 7.5" (16:9 aspect)
        # DPI = pixels / inches → 1980 / 10 = 198 DPI
        target_dpi = 198

        # Convert PDF to images at high resolution
        print(f"Converting to images at {target_dpi} DPI...")
        result = subprocess.run(
            ["pdftoppm", "-png", "-r", str(target_dpi), str(pdf_path), str(temp_dir / "slide")],
            capture_output=True,
            text=True,
        )
        if result.returncode != 0:
            raise RuntimeError("Image conversion failed")

        # Get visible slide images (pdftoppm skips hidden slides)
        visible_images = sorted(temp_dir.glob("slide-*.png"))

        # Map visible image index to actual slide index
        visible_to_actual = {}
        visible_idx = 0
        for actual_idx in range(total_slides):
            if actual_idx not in hidden_slides:
                visible_to_actual[actual_idx] = visible_idx
                visible_idx += 1

        # Process requested slides
        created_files = []
        for slide_idx in valid_indices:
            output_path = output_dir / f"slide-{slide_idx}.png"

            if slide_idx in hidden_slides:
                # Create placeholder for hidden slide
                placeholder_img = create_hidden_slide_placeholder((SINGLE_WIDTH, SINGLE_HEIGHT))
                placeholder_img.save(output_path, "PNG")
                created_files.append(str(output_path))
                print(f"  Created placeholder for hidden slide {slide_idx}")
            else:
                # Get the corresponding visible image
                vis_idx = visible_to_actual.get(slide_idx)
                if vis_idx is not None and vis_idx < len(visible_images):
                    with Image.open(visible_images[vis_idx]) as img:
                        # Resize to target dimensions
                        img_resized = img.resize(
                            (SINGLE_WIDTH, SINGLE_HEIGHT),
                            Image.Resampling.LANCZOS
                        )
                        img_resized.save(output_path, "PNG")
                        created_files.append(str(output_path))
                        print(f"  Created: {output_path}")
                else:
                    print(f"  Warning: Could not find image for slide {slide_idx}")

        # Print summary
        if len(created_files) == 1:
            print(f"Created single thumbnail:")
        else:
            print(f"Created {len(created_files)} thumbnail(s):")
        for f in created_files:
            print(f"  - {f}")


def get_placeholder_regions(pptx_path):
    """Extract ALL text regions from the presentation.

    Returns a tuple of (placeholder_regions, slide_dimensions).
    text_regions is a dict mapping slide indices to lists of text regions.
    Each region is a dict with 'left', 'top', 'width', 'height' in inches.
    slide_dimensions is a tuple of (width_inches, height_inches).
    """
    prs = Presentation(str(pptx_path))
    inventory = extract_text_inventory(pptx_path, prs)
    placeholder_regions = {}

    # Get actual slide dimensions in inches (EMU to inches conversion)
    slide_width_inches = (prs.slide_width or 9144000) / 914400.0
    slide_height_inches = (prs.slide_height or 5143500) / 914400.0

    for slide_key, shapes in inventory.items():
        # Extract slide index from "slide-N" format
        slide_idx = int(slide_key.split("-")[1])
        regions = []

        for shape_key, shape_data in shapes.items():
            # The inventory only contains shapes with text, so all shapes should be highlighted
            regions.append(
                {
                    "left": shape_data.left,
                    "top": shape_data.top,
                    "width": shape_data.width,
                    "height": shape_data.height,
                }
            )

        if regions:
            placeholder_regions[slide_idx] = regions

    return placeholder_regions, (slide_width_inches, slide_height_inches)


def convert_to_images(pptx_path, temp_dir, dpi, backend=None):
    """Convert PowerPoint to images via PDF, handling hidden slides.

    Args:
        pptx_path: Path to the PowerPoint file
        temp_dir: Temporary directory for intermediate files
        dpi: DPI for image conversion
        backend: 'libreoffice' | 'msoffice' | 'auto' | None
    """
    # Detect hidden slides
    print("Analyzing presentation...")
    prs = Presentation(str(pptx_path))
    total_slides = len(prs.slides)

    # Find hidden slides (1-based indexing for display)
    hidden_slides = {
        idx + 1
        for idx, slide in enumerate(prs.slides)
        if slide.element.get("show") == "0"
    }

    print(f"Total slides: {total_slides}")
    if hidden_slides:
        print(f"Hidden slides: {sorted(hidden_slides)}")

    # Convert to PDF using unified backend
    print("Converting to PDF...")
    pdf_path = convert_pptx_to_pdf(Path(pptx_path), temp_dir, backend)

    # Convert PDF to images
    print(f"Converting to images at {dpi} DPI...")
    result = subprocess.run(
        ["pdftoppm", "-jpeg", "-r", str(dpi), str(pdf_path), str(temp_dir / "slide")],
        capture_output=True,
        text=True,
    )
    if result.returncode != 0:
        raise RuntimeError("Image conversion failed")

    visible_images = sorted(temp_dir.glob("slide-*.jpg"))

    # Create full list with placeholders for hidden slides
    all_images = []
    visible_idx = 0

    # Get placeholder dimensions from first visible slide
    if visible_images:
        with Image.open(visible_images[0]) as img:
            placeholder_size = img.size
    else:
        placeholder_size = (1920, 1080)

    for slide_num in range(1, total_slides + 1):
        if slide_num in hidden_slides:
            # Create placeholder image for hidden slide
            placeholder_path = temp_dir / f"hidden-{slide_num:03d}.jpg"
            placeholder_img = create_hidden_slide_placeholder(placeholder_size)
            placeholder_img.save(placeholder_path, "JPEG")
            all_images.append(placeholder_path)
        else:
            # Use the actual visible slide image
            if visible_idx < len(visible_images):
                all_images.append(visible_images[visible_idx])
                visible_idx += 1

    return all_images


def create_grids(
    image_paths,
    cols,
    width,
    output_path,
    placeholder_regions=None,
    slide_dimensions=None,
):
    """Create multiple thumbnail grids from slide images, max cols×(cols+1) images per grid."""
    # Maximum images per grid is cols × (cols + 1) for better proportions
    max_images_per_grid = cols * (cols + 1)
    grid_files = []

    print(
        f"Creating grids with {cols} columns (max {max_images_per_grid} images per grid)"
    )

    # Split images into chunks
    for chunk_idx, start_idx in enumerate(
        range(0, len(image_paths), max_images_per_grid)
    ):
        end_idx = min(start_idx + max_images_per_grid, len(image_paths))
        chunk_images = image_paths[start_idx:end_idx]

        # Create grid for this chunk
        grid = create_grid(
            chunk_images, cols, width, start_idx, placeholder_regions, slide_dimensions
        )

        # Generate output filename
        if len(image_paths) <= max_images_per_grid:
            # Single grid - use base filename without suffix
            grid_filename = output_path
        else:
            # Multiple grids - insert index before extension with dash
            stem = output_path.stem
            suffix = output_path.suffix
            grid_filename = output_path.parent / f"{stem}-{chunk_idx + 1}{suffix}"

        # Save grid
        grid_filename.parent.mkdir(parents=True, exist_ok=True)
        grid.save(str(grid_filename), quality=JPEG_QUALITY)
        grid_files.append(str(grid_filename))

    return grid_files


def create_grid(
    image_paths,
    cols,
    width,
    start_slide_num=0,
    placeholder_regions=None,
    slide_dimensions=None,
):
    """Create thumbnail grid from slide images with optional placeholder outlining."""
    font_size = int(width * FONT_SIZE_RATIO)
    label_padding = int(font_size * LABEL_PADDING_RATIO)

    # Get dimensions
    with Image.open(image_paths[0]) as img:
        aspect = img.height / img.width
    height = int(width * aspect)

    # Calculate grid size
    rows = (len(image_paths) + cols - 1) // cols
    grid_w = cols * width + (cols + 1) * GRID_PADDING
    grid_h = rows * (height + font_size + label_padding * 2) + (rows + 1) * GRID_PADDING

    # Create grid
    grid = Image.new("RGB", (grid_w, grid_h), "white")
    draw = ImageDraw.Draw(grid)

    # Load font with size based on thumbnail width
    try:
        # Use Pillow's default font with size
        font = ImageFont.load_default(size=font_size)
    except Exception:
        # Fall back to basic default font if size parameter not supported
        font = ImageFont.load_default()

    # Place thumbnails
    for i, img_path in enumerate(image_paths):
        row, col = i // cols, i % cols
        x = col * width + (col + 1) * GRID_PADDING
        y_base = (
            row * (height + font_size + label_padding * 2) + (row + 1) * GRID_PADDING
        )

        # Add label with actual slide number
        label = f"{start_slide_num + i}"
        bbox = draw.textbbox((0, 0), label, font=font)
        text_w = bbox[2] - bbox[0]
        draw.text(
            (x + (width - text_w) // 2, y_base + label_padding),
            label,
            fill="black",
            font=font,
        )

        # Add thumbnail below label with proportional spacing
        y_thumbnail = y_base + label_padding + font_size + label_padding

        with Image.open(img_path) as img:
            # Get original dimensions before thumbnail
            orig_w, orig_h = img.size

            # Apply placeholder outlines if enabled
            if placeholder_regions and (start_slide_num + i) in placeholder_regions:
                # Convert to RGBA for transparency support
                if img.mode != "RGBA":
                    img = img.convert("RGBA")

                # Get the regions for this slide
                regions = placeholder_regions[start_slide_num + i]

                # Calculate scale factors using actual slide dimensions
                if slide_dimensions:
                    slide_width_inches, slide_height_inches = slide_dimensions
                else:
                    # Fallback: estimate from image size at CONVERSION_DPI
                    slide_width_inches = orig_w / CONVERSION_DPI
                    slide_height_inches = orig_h / CONVERSION_DPI

                x_scale = orig_w / slide_width_inches
                y_scale = orig_h / slide_height_inches

                # Create a highlight overlay
                overlay = Image.new("RGBA", img.size, (255, 255, 255, 0))
                overlay_draw = ImageDraw.Draw(overlay)

                # Highlight each placeholder region
                for region in regions:
                    # Convert from inches to pixels in the original image
                    px_left = int(region["left"] * x_scale)
                    px_top = int(region["top"] * y_scale)
                    px_width = int(region["width"] * x_scale)
                    px_height = int(region["height"] * y_scale)

                    # Draw highlight outline with red color and thick stroke
                    # Using a bright red outline instead of fill
                    stroke_width = max(
                        5, min(orig_w, orig_h) // 150
                    )  # Thicker proportional stroke width
                    overlay_draw.rectangle(
                        [(px_left, px_top), (px_left + px_width, px_top + px_height)],
                        outline=(255, 0, 0, 255),  # Bright red, fully opaque
                        width=stroke_width,
                    )

                # Composite the overlay onto the image using alpha blending
                img = Image.alpha_composite(img, overlay)
                # Convert back to RGB for JPEG saving
                img = img.convert("RGB")

            img.thumbnail((width, height), Image.Resampling.LANCZOS)
            w, h = img.size
            tx = x + (width - w) // 2
            ty = y_thumbnail + (height - h) // 2
            grid.paste(img, (tx, ty))

            # Add border
            if BORDER_WIDTH > 0:
                draw.rectangle(
                    [
                        (tx - BORDER_WIDTH, ty - BORDER_WIDTH),
                        (tx + w + BORDER_WIDTH - 1, ty + h + BORDER_WIDTH - 1),
                    ],
                    outline="gray",
                    width=BORDER_WIDTH,
                )

    return grid


if __name__ == "__main__":
    main()
