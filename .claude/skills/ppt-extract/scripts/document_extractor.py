#!/usr/bin/env python3
"""
문서 양식 추출기.

PPTX에서 문서 양식(슬라이드 마스터, 레이아웃, 테마)을 추출합니다.

PRD 요구사항:
- 표지, 목차, 내지(유형별 대표), 섹션 레이아웃 추출
- 슬라이드 마스터 + 레이아웃 OOXML 저장
- 테마 색상/폰트 추출
- 로고 등 미디어 에셋 추출
- 레이아웃 유형 분류 (LLM 판단)
"""

import json
import shutil
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any

try:
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.enum.shapes import PP_PLACEHOLDER
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

import sys
SCRIPT_DIR = Path(__file__).parent
sys.path.insert(0, str(SCRIPT_DIR.parent))

from shared.ooxml_parser import OOXMLParser, NAMESPACES
from shared.yaml_utils import save_yaml, generate_template_yaml
from scripts.thumbnail import convert_pptx_to_images, create_grid_view, THUMBNAIL_SIZES


@dataclass
class PlaceholderInfo:
    """플레이스홀더 정보."""
    id: str
    type: str
    idx: int
    position: Dict[str, float]


@dataclass
class ContentZone:
    """콘텐츠 영역."""
    x: float
    y: float
    width: float
    height: float


@dataclass
class LayoutInfo:
    """레이아웃 정보."""
    index: int
    name: str
    type: str  # cover, toc, section, body
    variant: Optional[str] = None  # title-only, title-action, two-column
    ooxml_file: str = ""
    thumbnail: str = ""
    placeholders: List[PlaceholderInfo] = field(default_factory=list)
    content_zone: Optional[ContentZone] = None

    def to_dict(self) -> Dict:
        """딕셔너리로 변환."""
        d = {
            'index': self.index,
            'name': self.name,
            'type': self.type,
        }
        if self.variant:
            d['variant'] = self.variant
        if self.ooxml_file:
            d['ooxml_file'] = self.ooxml_file
        if self.thumbnail:
            d['thumbnail'] = self.thumbnail
        if self.placeholders:
            d['placeholders'] = [
                {
                    'id': ph.id,
                    'type': ph.type,
                    'position': ph.position
                }
                for ph in self.placeholders
            ]
        if self.content_zone:
            d['content_zone'] = asdict(self.content_zone)
        return d


@dataclass
class CommonElement:
    """공통 요소 (로고 등)."""
    id: str
    type: str  # image, text
    file: Optional[str] = None
    position: Optional[Dict[str, float]] = None


@dataclass
class MasterInfo:
    """슬라이드 마스터 정보."""
    ooxml_file: str
    common_elements: List[CommonElement] = field(default_factory=list)

    def to_dict(self) -> Dict:
        """딕셔너리로 변환."""
        return {
            'ooxml_file': self.ooxml_file,
            'common_elements': [
                {
                    'id': ce.id,
                    'type': ce.type,
                    'file': ce.file,
                    'position': ce.position
                }
                for ce in self.common_elements
            ]
        }


@dataclass
class ThemeInfo:
    """테마 정보."""
    ooxml_file: str
    colors: Dict[str, str]
    fonts: Dict[str, str]

    def to_dict(self) -> Dict:
        """딕셔너리로 변환."""
        return {
            'ooxml_file': self.ooxml_file,
            'colors': self.colors,
            'fonts': self.fonts
        }


class DocumentExtractor:
    """문서 양식 추출기."""

    # 플레이스홀더 타입 매핑
    PH_TYPE_MAP = {
        1: 'title',
        2: 'body',
        3: 'ctrTitle',
        4: 'subTitle',
        5: 'dt',
        6: 'sldNum',
        7: 'ftr',
        8: 'hdr',
        9: 'obj',
        10: 'chart',
        11: 'tbl',
        12: 'clipArt',
        13: 'dgm',
        14: 'media',
        15: 'sldImg',
        16: 'pic',
    }

    def __init__(
        self,
        input_path: Path,
        group: str,
        name: str,
        output_path: Path,
        auto_classify: bool = False
    ):
        """
        Args:
            input_path: PPTX 파일 경로
            group: 문서 그룹명
            name: 템플릿 이름
            output_path: 출력 경로
            auto_classify: LLM 입력 없이 규칙 기반 분류 사용
        """
        if not HAS_PPTX:
            raise ImportError("python-pptx가 필요합니다: pip install python-pptx")

        self.input_path = Path(input_path)
        self.group = group
        self.name = name
        self.output_path = Path(output_path)
        self.auto_classify = auto_classify

        # OOXML 파서
        self.parser = OOXMLParser(self.input_path)

        # python-pptx Presentation
        self.prs = Presentation(str(self.input_path))

        # 슬라이드 크기
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height

    def run(self):
        """메인 추출 워크플로우."""
        print(f"문서 양식 추출 시작: {self.input_path.name}")
        print(f"출력 경로: {self.output_path}")

        # 1. 레이아웃 분석
        print("\n[1/7] 레이아웃 분석 중...")
        all_layouts = self.analyze_layouts()
        print(f"  발견된 레이아웃: {len(all_layouts)}개")

        # 2. 레이아웃 분류 (LLM 판단)
        print("\n[2/7] 레이아웃 분류 중 (LLM 판단)...")
        classified_layouts = self.classify_layouts_with_llm(all_layouts)

        # 3. 대표 레이아웃 선별
        print("\n[3/7] 대표 레이아웃 선별 중...")
        selected_layouts = self.select_representative_layouts(classified_layouts)
        print(f"  선별된 레이아웃: {len(selected_layouts)}개")
        for layout in selected_layouts:
            print(f"    - [{layout.index}] {layout.name}: {layout.type}" +
                  (f" ({layout.variant})" if layout.variant else ""))

        # 4. OOXML 추출
        print("\n[4/7] OOXML 추출 중...")
        self.extract_ooxml(selected_layouts)

        # 5. 에셋 추출
        print("\n[5/7] 에셋 추출 중...")
        master = self.extract_master_and_assets()

        # 6. 테마 추출
        print("\n[6/7] 테마 추출 중...")
        theme = self.extract_theme()

        # 7. 썸네일 생성
        print("\n[7/7] 썸네일 생성 중...")
        self.generate_thumbnail()

        # 메타데이터 저장
        print("\n메타데이터 저장 중...")
        self.save_metadata(selected_layouts, master, theme)

        # 레지스트리 업데이트
        print("\n레지스트리 업데이트 중...")
        try:
            from scripts.registry_manager import RegistryManager
            manager = RegistryManager()
            manager.update_document(self.output_path)
        except Exception as e:
            print(f"  Warning: 레지스트리 업데이트 실패: {e}")

        print(f"\n추출 완료!")

    def analyze_layouts(self) -> List[LayoutInfo]:
        """모든 레이아웃 분석."""
        layouts = []

        for i, layout in enumerate(self.prs.slide_layouts):
            # 플레이스홀더 추출
            placeholders = self._extract_placeholders(layout)

            layout_info = LayoutInfo(
                index=i,
                name=layout.name,
                type='unknown',  # 나중에 분류
                placeholders=placeholders
            )

            # 콘텐츠 영역 계산
            layout_info.content_zone = self._calculate_content_zone(placeholders)

            layouts.append(layout_info)

        return layouts

    def _extract_placeholders(self, layout) -> List[PlaceholderInfo]:
        """레이아웃에서 플레이스홀더 추출."""
        placeholders = []

        for shape in layout.placeholders:
            ph_format = shape.placeholder_format
            ph_type_int = ph_format.type if ph_format.type else 0

            # 타입 문자열 변환
            ph_type = self.PH_TYPE_MAP.get(int(ph_type_int), 'unknown')

            # 위치 계산 (퍼센트)
            position = {
                'x': round((shape.left / self.slide_width) * 100, 2),
                'y': round((shape.top / self.slide_height) * 100, 2),
                'width': round((shape.width / self.slide_width) * 100, 2),
                'height': round((shape.height / self.slide_height) * 100, 2),
            }

            placeholders.append(PlaceholderInfo(
                id=shape.name or f"ph_{ph_format.idx}",
                type=ph_type,
                idx=ph_format.idx,
                position=position
            ))

        return placeholders

    def _calculate_content_zone(self, placeholders: List[PlaceholderInfo]) -> ContentZone:
        """콘텐츠 영역 계산.

        PRD 규칙:
        - Title Zone: 상단 ~22%
        - Footer Zone: 하단 ~8%
        - Content Zone: Title 하단 + 2% ~ Footer 상단 - 2%
        """
        # 기본값
        title_bottom = 22.0
        footer_top = 92.0

        for ph in placeholders:
            pos = ph.position
            bottom = pos['y'] + pos['height']

            # 타이틀 영역 하단
            if ph.type in ['title', 'ctrTitle']:
                title_bottom = max(title_bottom, bottom)

            # 푸터 영역 상단
            if ph.type in ['ftr', 'sldNum', 'dt']:
                footer_top = min(footer_top, pos['y'])

        # 여백 추가 (2%)
        content_top = title_bottom + 2
        content_bottom = footer_top - 2

        return ContentZone(
            x=3.0,  # 좌우 여백 3%
            y=round(content_top, 2),
            width=94.0,
            height=round(content_bottom - content_top, 2)
        )

    def classify_layouts_with_llm(self, layouts: List[LayoutInfo]) -> List[LayoutInfo]:
        """LLM을 사용하여 레이아웃 분류.

        실제 LLM 호출은 CLI 환경에서 수행됨.
        여기서는 LLM에 전달할 데이터를 준비하고, 응답을 파싱합니다.

        auto_classify=True인 경우 규칙 기반 분류를 바로 사용합니다.
        """
        # 자동 모드: 규칙 기반 분류 바로 사용
        if self.auto_classify:
            print("  자동 모드: 규칙 기반 분류 사용")
            return self._classify_layouts_rule_based(layouts)

        # LLM 요청 데이터 준비
        layouts_data = []
        for layout in layouts:
            layouts_data.append({
                'index': layout.index,
                'name': layout.name,
                'placeholders': [
                    {'type': ph.type, 'position': ph.position}
                    for ph in layout.placeholders
                ],
                'placeholder_count': len(layout.placeholders)
            })

        # LLM 프롬프트 출력 (사용자/Claude가 판단)
        prompt = self._generate_classification_prompt(layouts_data)
        print("\n" + "="*60)
        print("LLM 레이아웃 분류 요청")
        print("="*60)
        print(prompt)
        print("="*60)

        # 사용자 입력 대기 (JSON 형식)
        print("\n레이아웃 분류 결과를 JSON으로 입력하세요:")
        print("형식: [{\"index\": 0, \"type\": \"cover\", \"variant\": null}, ...]")
        print("(또는 'auto'를 입력하면 규칙 기반 분류 사용)")

        try:
            user_input = input("\n> ").strip()

            if user_input.lower() == 'auto':
                # 규칙 기반 분류로 폴백
                return self._classify_layouts_rule_based(layouts)

            # JSON 파싱
            classifications = json.loads(user_input)

            # 분류 결과 적용
            for cls in classifications:
                idx = cls['index']
                for layout in layouts:
                    if layout.index == idx:
                        layout.type = cls['type']
                        layout.variant = cls.get('variant')
                        break

        except (json.JSONDecodeError, KeyError, EOFError) as e:
            print(f"Warning: 입력 처리 실패 ({e}), 규칙 기반 분류 사용")
            return self._classify_layouts_rule_based(layouts)

        return layouts

    def _generate_classification_prompt(self, layouts_data: List[Dict]) -> str:
        """LLM 분류 프롬프트 생성."""
        return f"""## 레이아웃 분류 요청

### 레이아웃 목록
{json.dumps(layouts_data, ensure_ascii=False, indent=2)}

### 분류 기준
- **cover**: 표지 (큰 제목 중앙 배치, ctrTitle 포함)
- **toc**: 목차 (번호가 있는 항목 목록)
- **section**: 섹션 구분 (제목만, 색상 배경)
- **body**: 내지 (콘텐츠 영역 존재)

### 내지(body) variant
- **title-only**: 제목만
- **title-action**: 제목 + 액션타이틀 (상단 본문)
- **title-body**: 제목 + 본문
- **two-column**: 2단 레이아웃
- **image-text**: 이미지 + 텍스트

### 출력 형식 (JSON 배열)
[
  {{"index": 0, "type": "cover", "variant": null}},
  {{"index": 1, "type": "toc", "variant": null}},
  {{"index": 2, "type": "section", "variant": null}},
  {{"index": 3, "type": "body", "variant": "title-action"}},
  ...
]"""

    def _classify_layouts_rule_based(self, layouts: List[LayoutInfo]) -> List[LayoutInfo]:
        """규칙 기반 레이아웃 분류 (폴백)."""
        for layout in layouts:
            layout.type, layout.variant = self._classify_single_layout(layout)
        return layouts

    def _classify_single_layout(self, layout: LayoutInfo) -> Tuple[str, Optional[str]]:
        """단일 레이아웃 규칙 기반 분류."""
        name_lower = layout.name.lower()
        ph_types = [ph.type for ph in layout.placeholders]

        # 1. 표지 판별
        if 'ctrTitle' in ph_types:
            return ('cover', None)
        if any(kw in name_lower for kw in ['표지', 'cover', 'title slide', '타이틀']):
            return ('cover', None)

        # 2. 목차 판별 (간지 = 목차/섹션 구분 페이지)
        if any(kw in name_lower for kw in ['목차', 'toc', 'contents', 'agenda', '개요', '간지']):
            return ('toc', None)

        # 3. 섹션 판별
        if any(kw in name_lower for kw in ['섹션', 'section', 'divider', '구분']):
            return ('section', None)

        # 4. 내지 유형 분류
        return self._classify_body_variant(layout)

    def _classify_body_variant(self, layout: LayoutInfo) -> Tuple[str, str]:
        """내지 세부 유형 분류."""
        phs = layout.placeholders
        ph_count = len(phs)

        has_title = any(ph.type in ['title', 'ctrTitle'] for ph in phs)
        body_phs = [ph for ph in phs if ph.type == 'body']

        # 제목만
        if ph_count == 1 and has_title:
            return ('body', 'title-only')

        # 2단 레이아웃
        if len(body_phs) >= 2:
            return ('body', 'two-column')

        # 이미지 + 텍스트
        if any(ph.type == 'pic' for ph in phs):
            return ('body', 'image-text')

        # 제목 + 액션타이틀 (본문이 상단에 있는 경우)
        if has_title and body_phs:
            body_y = body_phs[0].position['y']
            if body_y < 30:  # 상단 30% 이내
                return ('body', 'title-action')
            return ('body', 'title-body')

        return ('body', 'generic')

    def select_representative_layouts(self, layouts: List[LayoutInfo]) -> List[LayoutInfo]:
        """유형별 대표 레이아웃 선별.

        PRD 규칙:
        - 표지: 최대 2장
        - 목차: 1장
        - 섹션: 1장
        - 내지: variant별 1장씩 (총 3~5장)
        """
        selected = []

        # 유형별 분류
        covers = [l for l in layouts if l.type == 'cover']
        tocs = [l for l in layouts if l.type == 'toc']
        sections = [l for l in layouts if l.type == 'section']
        bodies = [l for l in layouts if l.type == 'body']

        # 표지 선택 (최대 2장)
        selected.extend(covers[:2])

        # 목차 선택 (1장)
        if tocs:
            selected.append(tocs[0])

        # 섹션 선택 (1장)
        if sections:
            selected.append(sections[0])

        # 내지 선택 (variant별 1장)
        seen_variants = set()
        for body in bodies:
            variant = body.variant or 'generic'
            if variant not in seen_variants:
                selected.append(body)
                seen_variants.add(variant)

        return selected

    def extract_ooxml(self, layouts: List[LayoutInfo]):
        """선별된 레이아웃의 OOXML 추출."""
        ooxml_dir = self.output_path / 'ooxml'
        ooxml_dir.mkdir(parents=True, exist_ok=True)
        rels_dir = ooxml_dir / '_rels'
        rels_dir.mkdir(parents=True, exist_ok=True)

        for layout in layouts:
            # 레이아웃 XML (1-based 인덱스)
            layout_idx = layout.index + 1
            layout_path = f'ppt/slideLayouts/slideLayout{layout_idx}.xml'
            rels_path = f'ppt/slideLayouts/_rels/slideLayout{layout_idx}.xml.rels'

            # XML 추출
            xml_content = self.parser.read_xml_string(layout_path, pretty=True)
            if xml_content:
                output_file = ooxml_dir / f'slideLayout{layout_idx}.xml'
                output_file.write_text(xml_content, encoding='utf-8')
                layout.ooxml_file = f'ooxml/slideLayout{layout_idx}.xml'
                print(f"  추출: {output_file.name}")

            # 관계 파일 추출
            rels_content = self.parser.read_xml_string(rels_path, pretty=True)
            if rels_content:
                rels_output = rels_dir / f'slideLayout{layout_idx}.xml.rels'
                rels_output.write_text(rels_content, encoding='utf-8')

        # 슬라이드 마스터 추출
        master_xml = self.parser.read_xml_string('ppt/slideMasters/slideMaster1.xml', pretty=True)
        if master_xml:
            (ooxml_dir / 'slideMaster1.xml').write_text(master_xml, encoding='utf-8')
            print(f"  추출: slideMaster1.xml")

        master_rels = self.parser.read_xml_string('ppt/slideMasters/_rels/slideMaster1.xml.rels', pretty=True)
        if master_rels:
            (rels_dir / 'slideMaster1.xml.rels').write_text(master_rels, encoding='utf-8')

        # 테마 추출
        theme_xml = self.parser.read_xml_string('ppt/theme/theme1.xml', pretty=True)
        if theme_xml:
            (ooxml_dir / 'theme1.xml').write_text(theme_xml, encoding='utf-8')
            print(f"  추출: theme1.xml")

    def extract_master_and_assets(self) -> MasterInfo:
        """슬라이드 마스터 및 모든 미디어 에셋 추출."""
        assets_dir = self.output_path / 'assets' / 'media'
        assets_dir.mkdir(parents=True, exist_ok=True)

        common_elements = []

        # 1. PPTX 내 모든 미디어 파일 추출 (ppt/media/*)
        import zipfile
        with zipfile.ZipFile(self.input_path, 'r') as zf:
            media_files = [f for f in zf.namelist() if f.startswith('ppt/media/')]

            for media_path in media_files:
                filename = Path(media_path).name
                data = zf.read(media_path)
                output_file = assets_dir / filename
                output_file.write_bytes(data)
                print(f"  추출: {filename} ({len(data):,} bytes)")

                # 이미지 타입 판별
                ext = Path(filename).suffix.lower()
                media_type = 'image' if ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.emf', '.wmf', '.svg'] else 'media'

                common_elements.append(CommonElement(
                    id=filename,
                    type=media_type,
                    file=f"assets/media/{filename}"
                ))

        # 2. 마스터에서 사용하는 미디어 참조 정보 (위치 정보 등)
        refs = self.parser.parse_rels('ppt/slideMasters/_rels/slideMaster1.xml.rels')
        master_refs = {ref.rid: ref for ref in refs if ref.type == 'image'}

        # common_elements에 마스터 참조 ID 업데이트
        for elem in common_elements:
            for rid, ref in master_refs.items():
                if elem.file.endswith(Path(ref.target).name):
                    elem.id = rid
                    break

        return MasterInfo(
            ooxml_file='ooxml/slideMaster1.xml',
            common_elements=common_elements
        )

    def extract_theme(self) -> ThemeInfo:
        """테마 색상/폰트 추출."""
        colors = self.parser.extract_theme_colors()
        fonts = self.parser.extract_theme_fonts()

        print(f"  색상: {len(colors)}개")
        print(f"  폰트: {fonts.get('major', 'N/A')} / {fonts.get('minor', 'N/A')}")

        return ThemeInfo(
            ooxml_file='ooxml/theme1.xml',
            colors=colors,
            fonts=fonts
        )

    def save_metadata(
        self,
        layouts: List[LayoutInfo],
        master: MasterInfo,
        theme: ThemeInfo
    ):
        """template.yaml 저장."""
        # 슬라이드 크기 정보
        width_px, height_px = self.parser.get_slide_size_px()
        slide_size = {
            'width': width_px,
            'height': height_px,
            'ratio': self.parser.get_aspect_ratio()
        }

        # 템플릿 데이터 생성
        template_data = generate_template_yaml(
            doc_id=f"{self.group}-{self.name}",
            doc_name=self.name,
            group=self.group,
            source_file=self.input_path.name,
            slide_size=slide_size,
            layouts=[l.to_dict() for l in layouts],
            master=master.to_dict(),
            theme=theme.to_dict()
        )

        # YAML 저장
        yaml_path = self.output_path / 'template.yaml'
        save_yaml(
            template_data,
            yaml_path,
            header=f"문서 양식: {self.name}\n추출 원본: {self.input_path.name}"
        )
        print(f"  저장: {yaml_path}")

    def generate_thumbnail(self):
        """문서 양식 썸네일 생성 (그리드 뷰).

        PPTX의 모든 슬라이드를 그리드 형태로 배치한 썸네일 생성.
        LibreOffice + pdftoppm 필요.
        """
        import tempfile

        thumbnail_path = self.output_path / 'thumbnail.jpg'

        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_path = Path(temp_dir)

                # PPTX → 이미지 변환
                images = convert_pptx_to_images(self.input_path, temp_path)
                print(f"  슬라이드 {len(images)}개 변환됨")

                if images:
                    # 그리드 뷰 생성 (최대 5열)
                    cols = min(5, len(images))
                    create_grid_view(images, thumbnail_path, cols=cols)
                    print(f"  저장: {thumbnail_path.name}")

        except FileNotFoundError as e:
            print(f"  경고: LibreOffice/pdftoppm 미설치 - 썸네일 생성 생략")
            print(f"    설치: sudo apt install libreoffice poppler-utils")
        except Exception as e:
            print(f"  경고: 썸네일 생성 실패 - {e}")


def main():
    """CLI 엔트리포인트."""
    import argparse

    parser = argparse.ArgumentParser(
        description='문서 양식 추출기 - PPTX에서 문서 양식(슬라이드 마스터, 레이아웃, 테마)을 추출합니다.',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
예시:
  python document_extractor.py input.pptx --group 동국그룹 --name 기본양식 --output templates/documents/dongkuk/standard
  python document_extractor.py input.pptx --group default --name my-template --output output/ --auto-classify
        """
    )
    parser.add_argument('input', help='PPTX 파일 경로')
    parser.add_argument('--group', required=True, help='문서 그룹명 (예: 동국그룹, default)')
    parser.add_argument('--name', required=True, help='템플릿 이름 (예: 기본양식, kkalkkumi-deepgreen)')
    parser.add_argument('--output', required=True, help='출력 디렉토리 경로')
    parser.add_argument('--auto-classify', action='store_true',
                       help='LLM 입력 없이 규칙 기반 분류 사용')

    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"오류: 입력 파일을 찾을 수 없습니다: {input_path}")
        sys.exit(1)

    extractor = DocumentExtractor(
        input_path=input_path,
        group=args.group,
        name=args.name,
        output_path=Path(args.output),
        auto_classify=args.auto_classify
    )
    extractor.run()


if __name__ == "__main__":
    main()
