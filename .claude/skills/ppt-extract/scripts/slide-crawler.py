#!/usr/bin/env python3
"""
PPTX 슬라이드 파싱 및 도형/텍스트 추출.

슬라이드에서 모든 도형(shape)을 추출하고 구조화된 JSON으로 출력.
GroupShape 재귀 처리, 절대 위치 계산, 영역 분류(Title/Content/Footer) 포함.

Usage:
    python slide-crawler.py input.pptx --slide 3 --output working/parsed.json
    python slide-crawler.py input.pptx --all --output working/all-slides.json
"""

import argparse
import json
import sys
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER


# EMU to inches conversion factor
EMU_PER_INCH = 914400


@dataclass
class TextRun:
    """텍스트 런 (폰트 스타일 포함)"""
    text: str
    font_name: Optional[str] = None
    font_size: Optional[float] = None  # points
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    color: Optional[str] = None  # RGB hex


@dataclass
class Paragraph:
    """단락 정보"""
    text: str
    runs: List[TextRun] = field(default_factory=list)
    alignment: Optional[str] = None
    level: int = 0
    bullet: bool = False
    space_before: Optional[float] = None  # points
    space_after: Optional[float] = None  # points
    line_spacing: Optional[float] = None  # points


@dataclass
class ShapeGeometry:
    """도형 위치/크기 (vmin 및 EMU 단위)"""
    x: float  # vmin
    y: float  # vmin
    cx: float  # vmin (width)
    cy: float  # vmin (height)
    emu: Dict[str, int] = field(default_factory=dict)  # 원본 EMU 값


@dataclass
class ShapeStyle:
    """도형 스타일 정보"""
    fill_color: Optional[str] = None  # RGB hex
    fill_type: Optional[str] = None  # solid, gradient, picture, none
    line_color: Optional[str] = None
    line_width: Optional[float] = None  # points
    shadow: bool = False
    rotation: float = 0.0  # degrees


@dataclass
class ExtractedShape:
    """추출된 도형 정보"""
    id: str
    name: str
    shape_type: str
    geometry: ShapeGeometry
    style: ShapeStyle
    zone: str  # title, content, footer
    paragraphs: List[Paragraph] = field(default_factory=list)
    placeholder_type: Optional[str] = None
    is_group: bool = False
    children: List['ExtractedShape'] = field(default_factory=list)


@dataclass
class ExtractedSlide:
    """추출된 슬라이드 정보"""
    index: int
    width: float  # vmin
    height: float  # vmin
    width_emu: int
    height_emu: int
    shapes: List[ExtractedShape] = field(default_factory=list)
    content_zone: Optional[Dict[str, float]] = None  # top, bottom in vmin


def emu_to_inches(emu: int) -> float:
    """EMU to inches"""
    return emu / EMU_PER_INCH


def emu_to_vmin(emu: int, vmin_emu: int) -> float:
    """EMU to vmin (relative to slide's shorter dimension)"""
    return round((emu / vmin_emu) * 100, 2)


def get_rgb_color(color_format) -> Optional[str]:
    """pptx 컬러 객체에서 RGB hex 추출"""
    try:
        if color_format and hasattr(color_format, 'rgb') and color_format.rgb:
            return str(color_format.rgb)
    except Exception:
        pass
    return None


def get_placeholder_type_name(placeholder_format) -> Optional[str]:
    """플레이스홀더 타입 이름 반환"""
    if not placeholder_format or not placeholder_format.type:
        return None

    type_str = str(placeholder_format.type)
    # "TITLE (1)" -> "TITLE"
    return type_str.split('(')[0].strip().split('.')[-1]


def determine_zone(shape, slide_height_emu: int) -> str:
    """도형이 속한 영역 판별 (title/content/footer)"""

    # 플레이스홀더 기반 판별
    if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
        ph_type = shape.placeholder_format.type
        if ph_type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.SUBTITLE]:
            return "title"
        if ph_type in [PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.SLIDE_NUMBER, PP_PLACEHOLDER.DATE]:
            return "footer"

    # 이름 기반 판별
    name_lower = shape.name.lower() if hasattr(shape, 'name') else ""
    if 'title' in name_lower:
        return "title"
    if 'footer' in name_lower or 'slide number' in name_lower:
        return "footer"

    # 위치 기반 판별
    if hasattr(shape, 'top'):
        shape_center_y = shape.top + (shape.height / 2) if hasattr(shape, 'height') else shape.top
        relative_y = shape_center_y / slide_height_emu

        if relative_y < 0.22:
            return "title"
        if relative_y > 0.90:
            return "footer"

    return "content"


def extract_text_runs(paragraph) -> List[TextRun]:
    """단락에서 텍스트 런 추출"""
    runs = []
    for run in paragraph.runs:
        font = run.font
        runs.append(TextRun(
            text=run.text,
            font_name=font.name if font.name else None,
            font_size=font.size.pt if font.size else None,
            bold=font.bold,
            italic=font.italic,
            color=get_rgb_color(font.color) if hasattr(font, 'color') else None
        ))
    return runs


def extract_paragraph(paragraph) -> Paragraph:
    """단락 정보 추출"""
    alignment_map = {
        1: "LEFT",
        2: "CENTER",
        3: "RIGHT",
        4: "JUSTIFY"
    }

    # 불릿 감지
    bullet = False
    if hasattr(paragraph, '_p') and paragraph._p is not None:
        pPr = paragraph._p.pPr
        if pPr is not None:
            ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            if pPr.find(f"{ns}buChar") is not None or pPr.find(f"{ns}buAutoNum") is not None:
                bullet = True

    return Paragraph(
        text=paragraph.text.strip(),
        runs=extract_text_runs(paragraph),
        alignment=alignment_map.get(paragraph.alignment, None) if paragraph.alignment else None,
        level=paragraph.level if hasattr(paragraph, 'level') else 0,
        bullet=bullet,
        space_before=paragraph.space_before.pt if paragraph.space_before else None,
        space_after=paragraph.space_after.pt if paragraph.space_after else None,
        line_spacing=paragraph.line_spacing.pt if hasattr(paragraph, 'line_spacing') and paragraph.line_spacing else None
    )


def extract_shape_style(shape) -> ShapeStyle:
    """도형 스타일 추출"""
    style = ShapeStyle()

    try:
        if hasattr(shape, 'fill'):
            fill = shape.fill
            if fill.type is not None:
                style.fill_type = str(fill.type).split('.')[-1].lower()
            if hasattr(fill, 'fore_color') and fill.fore_color:
                style.fill_color = get_rgb_color(fill.fore_color)
    except Exception:
        pass

    try:
        if hasattr(shape, 'line'):
            line = shape.line
            if line.color and line.color.rgb:
                style.line_color = str(line.color.rgb)
            if line.width:
                style.line_width = line.width.pt
    except Exception:
        pass

    try:
        if hasattr(shape, 'shadow') and shape.shadow:
            style.shadow = shape.shadow.inherit
    except Exception:
        pass

    try:
        if hasattr(shape, 'rotation'):
            style.rotation = shape.rotation
    except Exception:
        pass

    return style


def extract_shape(
    shape,
    vmin_emu: int,
    slide_height_emu: int,
    parent_left: int = 0,
    parent_top: int = 0,
    shape_counter: List[int] = None
) -> ExtractedShape:
    """도형 정보 추출 (재귀적으로 GroupShape 처리)"""

    if shape_counter is None:
        shape_counter = [0]

    shape_id = f"shape-{shape_counter[0]}"
    shape_counter[0] += 1

    # 절대 위치 계산
    left_emu = (shape.left if hasattr(shape, 'left') else 0) + parent_left
    top_emu = (shape.top if hasattr(shape, 'top') else 0) + parent_top
    width_emu = shape.width if hasattr(shape, 'width') else 0
    height_emu = shape.height if hasattr(shape, 'height') else 0

    geometry = ShapeGeometry(
        x=emu_to_vmin(left_emu, vmin_emu),
        y=emu_to_vmin(top_emu, vmin_emu),
        cx=emu_to_vmin(width_emu, vmin_emu),
        cy=emu_to_vmin(height_emu, vmin_emu),
        emu={"x": left_emu, "y": top_emu, "cx": width_emu, "cy": height_emu}
    )

    # shape type 추출
    shape_type = "unknown"
    if hasattr(shape, 'shape_type') and shape.shape_type:
        shape_type = str(shape.shape_type).split('.')[-1]

    # 플레이스홀더 타입
    placeholder_type = None
    if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
        placeholder_type = get_placeholder_type_name(shape.placeholder_format)

    # 텍스트 추출
    paragraphs = []
    if hasattr(shape, 'text_frame') and shape.text_frame:
        for para in shape.text_frame.paragraphs:
            if para.text.strip():
                paragraphs.append(extract_paragraph(para))

    # GroupShape 처리
    children = []
    is_group = False
    if hasattr(shape, 'shapes'):
        is_group = True
        for child in shape.shapes:
            children.append(extract_shape(
                child, vmin_emu, slide_height_emu,
                parent_left=left_emu, parent_top=top_emu,
                shape_counter=shape_counter
            ))

    return ExtractedShape(
        id=shape_id,
        name=shape.name if hasattr(shape, 'name') else "",
        shape_type=shape_type,
        geometry=geometry,
        style=extract_shape_style(shape),
        zone=determine_zone(shape, slide_height_emu),
        paragraphs=paragraphs,
        placeholder_type=placeholder_type,
        is_group=is_group,
        children=children
    )


def detect_content_zone(shapes: List[ExtractedShape], slide_height_vmin: float) -> Dict[str, float]:
    """콘텐츠 영역 경계 계산 (vmin 단위)"""

    title_shapes = [s for s in shapes if s.zone == "title"]
    footer_shapes = [s for s in shapes if s.zone == "footer"]

    # Title 영역 하단
    if title_shapes:
        title_bottom = max(s.geometry.y + s.geometry.cy for s in title_shapes)
        content_top = title_bottom + 2.0  # 2% 여백
    else:
        content_top = 22.0  # 기본 22%

    # Footer 영역 상단
    if footer_shapes:
        footer_top = min(s.geometry.y for s in footer_shapes)
        content_bottom = footer_top - 2.0  # 2% 여백
    else:
        content_bottom = 92.0  # 기본 92%

    return {
        "top": round(content_top, 2),
        "bottom": round(content_bottom, 2)
    }


def extract_slide(slide, index: int) -> ExtractedSlide:
    """슬라이드 전체 추출"""

    # 슬라이드 크기
    prs = slide.part.package.presentation_part.presentation
    width_emu = prs.slide_width
    height_emu = prs.slide_height
    vmin_emu = min(width_emu, height_emu)

    # vmin 단위로 변환
    width_vmin = emu_to_vmin(width_emu, vmin_emu)
    height_vmin = emu_to_vmin(height_emu, vmin_emu)

    # 도형 추출
    shape_counter = [0]
    shapes = []
    for shape in slide.shapes:
        shapes.append(extract_shape(
            shape, vmin_emu, height_emu,
            shape_counter=shape_counter
        ))

    # 콘텐츠 영역 계산
    content_zone = detect_content_zone(shapes, height_vmin)

    return ExtractedSlide(
        index=index,
        width=width_vmin,
        height=height_vmin,
        width_emu=width_emu,
        height_emu=height_emu,
        shapes=shapes,
        content_zone=content_zone
    )


def dataclass_to_dict(obj) -> Any:
    """dataclass를 dict로 변환 (None 값 제외)"""
    if hasattr(obj, '__dataclass_fields__'):
        result = {}
        for key, value in asdict(obj).items():
            clean_value = dataclass_to_dict(value)
            if clean_value is not None:
                result[key] = clean_value
        return result
    elif isinstance(obj, list):
        return [dataclass_to_dict(item) for item in obj if dataclass_to_dict(item) is not None]
    elif isinstance(obj, dict):
        return {k: dataclass_to_dict(v) for k, v in obj.items() if dataclass_to_dict(v) is not None}
    else:
        return obj


def main():
    parser = argparse.ArgumentParser(
        description="PPTX 슬라이드 파싱 및 도형/텍스트 추출"
    )
    parser.add_argument("input", help="입력 PPTX 파일")
    parser.add_argument("--slide", type=int, help="특정 슬라이드 번호 (0-based)")
    parser.add_argument("--all", action="store_true", help="모든 슬라이드 추출")
    parser.add_argument("--output", "-o", help="출력 JSON 파일")
    parser.add_argument("--content-only", action="store_true", help="콘텐츠 영역만 추출")

    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: 파일을 찾을 수 없습니다: {args.input}")
        sys.exit(1)

    if not args.slide and not args.all:
        print("Error: --slide 또는 --all 옵션을 지정하세요")
        sys.exit(1)

    try:
        prs = Presentation(str(input_path))

        slides_data = []

        if args.all:
            for idx, slide in enumerate(prs.slides):
                extracted = extract_slide(slide, idx)
                slides_data.append(extracted)
        else:
            if args.slide < 0 or args.slide >= len(prs.slides):
                print(f"Error: 슬라이드 {args.slide}가 존재하지 않습니다 (0-{len(prs.slides)-1})")
                sys.exit(1)

            slide = prs.slides[args.slide]
            extracted = extract_slide(slide, args.slide)
            slides_data.append(extracted)

        # content-only 필터링
        if args.content_only:
            for slide_data in slides_data:
                slide_data.shapes = [
                    s for s in slide_data.shapes
                    if s.zone == "content"
                ]

        # 출력
        output_data = {
            "source_file": str(input_path.name),
            "slides": [dataclass_to_dict(s) for s in slides_data]
        }

        if args.output:
            output_path = Path(args.output)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(output_data, f, indent=2, ensure_ascii=False)
            print(f"저장됨: {args.output}")
        else:
            print(json.dumps(output_data, indent=2, ensure_ascii=False))

        # 통계
        total_shapes = sum(len(s.shapes) for s in slides_data)
        print(f"추출 완료: {len(slides_data)} 슬라이드, {total_shapes} 도형")

    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
