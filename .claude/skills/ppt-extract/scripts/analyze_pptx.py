#!/usr/bin/env python3
"""
PPTX 콘텐츠 추출 및 분석 도구

슬라이드를 분석하여 콘텐츠 카테고리로 분류하고 필요한 정보를 추출합니다.
"""

import json
import shutil
import sys
from pathlib import Path
from typing import Dict, List, Any
from zipfile import ZipFile

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Error: python-pptx가 필요합니다: pip install python-pptx")
    sys.exit(1)


def extract_slide_ooxml(pptx_path: Path, slide_index: int, output_dir: Path) -> Path:
    """
    특정 슬라이드의 OOXML을 추출하여 별도 파일로 저장.

    Args:
        pptx_path: PPTX 파일 경로
        slide_index: 슬라이드 인덱스 (0-based)
        output_dir: 출력 디렉토리

    Returns:
        추출된 OOXML 폴더 경로
    """
    # PPTX는 ZIP 파일
    with ZipFile(pptx_path, 'r') as zip_ref:
        # 슬라이드 XML 경로: ppt/slides/slide{N}.xml
        slide_xml_path = f"ppt/slides/slide{slide_index + 1}.xml"

        if slide_xml_path not in [n for n in zip_ref.namelist() if n.startswith('ppt/slides/')]:
            raise ValueError(f"슬라이드 {slide_index}를 찾을 수 없습니다")

        # 출력 폴더 생성
        output_dir.mkdir(parents=True, exist_ok=True)

        # 슬라이드 XML 추출
        xml_data = zip_ref.read(slide_xml_path)
        (output_dir / "slide.xml").write_bytes(xml_data)

        # 관련 레이아웃/마스터도 추출
        # 슬라이드 레이아웃 참조 찾기
        import xml.etree.ElementTree as ET
        root = ET.fromstring(xml_data)

        # 네임스페이스
        ns = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        }

        # 레이아웃 참조 추출
        for elem in root.findall('.//p:sldLayoutId', ns):
            # 관계 파일에서 레이아웃 ID 찾기
            pass

        # 레이아웃 XML 추출 (필요시)
        # 마스터 XML 추출 (필요시)

        return output_dir


def analyze_slide_content(slide) -> Dict[str, Any]:
    """
    슬라이드 콘텐츠를 분석하여 카테고리를 결정하고 메타데이터 추출.

    Returns:
        분석 결과 딕셔너리
    """
    info = {
        'shapes_count': len(slide.shapes),
        'has_title': False,
        'has_body': False,
        'has_image': False,
        'has_table': False,
        'has_chart': False,
        'has_group': False,
        'text_shapes': 0,
        'image_shapes': 0,
        'table_shapes': 0,
        'chart_shapes': 0,
        'group_shapes': 0,
        'title_text': '',
        'body_texts': [],
        'layout_name': slide.slide_layout.name if slide.slide_layout else 'Unknown'
    }

    # 도형 분석
    for shape in slide.shapes:
        # 그룹 도형
        if shape.shape_type == 6:  # Group
            info['has_group'] = True
            info['group_shapes'] += 1

        # 텍스트 프레임이 있는 도형
        if hasattr(shape, "text") and shape.text.strip():
            info['text_shapes'] += 1

            # 제목/부제목 확인
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                placeholder_type = shape.placeholder_format.type
                if placeholder_type == 0:  # Title
                    info['has_title'] = True
                    info['title_text'] = shape.text.strip()
                elif placeholder_type == 1:  # Centered Title
                    info['has_title'] = True
                    info['title_text'] = shape.text.strip()
                elif placeholder_type in [2, 3]:  # Body
                    info['has_body'] = True
                    info['body_texts'].append(shape.text.strip())
            else:
                # 일반 텍스트
                if not info['title_text'] and len(shape.text.strip()) < 100:
                    info['title_text'] = shape.text.strip()
                else:
                    info['has_body'] = True
                    info['body_texts'].append(shape.text.strip())

        # 이미지
        if shape.shape_type == 13:  # Picture
            info['has_image'] = True
            info['image_shapes'] += 1

        # 표
        if shape.shape_type == 19:  # Table
            info['has_table'] = True
            info['table_shapes'] += 1

        # 차트
        if shape.shape_type == 3:  # Chart
            info['has_chart'] = True
            info['chart_shapes'] += 1

    return info


def classify_slide(info: Dict[str, Any]) -> str:
    """
    슬라이드 분석 정보를 기반으로 카테고리 분류.

    Returns:
        카테고리 이름 (cover, toc, section-title, text, image, chart, diagram, table, etc.)
    """
    # 빈 슬라이드
    if info['shapes_count'] == 0:
        return 'empty'

    # 표지 (Cover)
    if info['has_title'] and not info['has_body'] and info['shapes_count'] <= 3:
        if info['has_image'] or info['title_text']:
            return 'cover'

    # 섹션 타이틀
    if info['has_title'] and not info['has_body'] and info['shapes_count'] <= 5:
        # 섹션 구분용 슬라이드
        return 'section-title'

    # 차트 슬라이드
    if info['has_chart']:
        return 'chart'

    # 표 슬라이드
    if info['has_table']:
        return 'table'

    # 이미지 메인 슬라이드
    if info['has_image'] and info['image_shapes'] >= 1 and not info['has_body']:
        return 'image'

    # 다이어그램 (그룹 도형이 많은 경우)
    if info['has_group'] and info['group_shapes'] >= 1:
        return 'diagram'

    # TOC (목차)
    title_lower = info['title_text'].lower()
    if '목차' in title_lower or 'toc' in title_lower or 'index' in title_lower or 'content' in title_lower:
        return 'toc'

    # 기본 텍스트 슬라이드
    if info['has_body']:
        # 불릿 포인트가 많은지 확인
        return 'text'

    # 기본값
    return 'text'


def categorize_all_slides(pptx_path: Path) -> List[Dict[str, Any]]:
    """
    모든 슬라이드를 분석하고 카테고리 분류.

    Returns:
        슬라이드 분석 결과 리스트
    """
    prs = Presentation(pptx_path)
    results = []

    print(f"총 {len(prs.slides)}개 슬라이드 분석 중...\n")

    for i, slide in enumerate(prs.slides):
        info = analyze_slide_content(slide)
        category = classify_slide(info)

        result = {
            'slide_index': i,
            'category': category,
            'layout': info['layout_name'],
            'shapes_count': info['shapes_count'],
            'has_title': info['has_title'],
            'title_text': info['title_text'],
            'has_body': info['has_body'],
            'has_image': info['has_image'],
            'has_table': info['has_table'],
            'has_chart': info['has_chart'],
            'has_group': info['has_group'],
        }

        results.append(result)

        # 요약 출력
        category_marker = {
            'cover': '📘',
            'toc': '📋',
            'section-title': '📌',
            'text': '📝',
            'image': '🖼️',
            'chart': '📊',
            'table': '📋',
            'diagram': '🔷',
            'empty': '⬜'
        }.get(category, '📄')

        print(f"{category_marker} 슬라이드 {i:2d}: {category:15s} | {info['layout']:30s} | 도형:{info['shapes_count']:2d} | {info['title_text'][:40]}")

    return results


def group_by_category(results: List[Dict[str, Any]]) -> Dict[str, List[Dict[str, Any]]]:
    """카테고리별로 슬라이드 그룹화"""
    grouped = {}
    for result in results:
        cat = result['category']
        if cat not in grouped:
            grouped[cat] = []
        grouped[cat].append(result)
    return grouped


def main():
    import argparse

    parser = argparse.ArgumentParser(description="PPTX 콘텐츠 분석 및 추출")
    parser.add_argument("input", help="입력 PPTX 파일")
    parser.add_argument("--output", "-o", help="출력 JSON 파일")
    parser.add_argument("--extract-ooxml", action="store_true", help="OOXML 추출")

    args = parser.parse_args()

    pptx_path = Path(args.input)
    if not pptx_path.exists():
        print(f"Error: 파일을 찾을 수 없습니다: {args.input}")
        sys.exit(1)

    # 슬라이드 분석
    results = categorize_all_slides(pptx_path)

    # 카테고리별 요약
    print("\n" + "=" * 80)
    print("카테고리별 요약")
    print("=" * 80)

    grouped = group_by_category(results)
    for cat, slides in sorted(grouped.items()):
        print(f"{cat:15s}: {len(slides):2d}개 슬라이드 (슬라이드 {[s['slide_index'] for s in slides]})")

    # JSON 저장
    if args.output:
        output_path = Path(args.output)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        summary = {
            'source_file': str(pptx_path),
            'total_slides': len(results),
            'categories': {cat: {'count': len(slides), 'slides': [s['slide_index'] for s in slides]}
                          for cat, slides in sorted(grouped.items())},
            'slides': results
        }

        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(summary, f, ensure_ascii=False, indent=2)

        print(f"\n저장됨: {output_path}")


if __name__ == "__main__":
    main()
