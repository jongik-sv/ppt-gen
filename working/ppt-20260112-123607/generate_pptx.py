#!/usr/bin/env python3
"""
HTML slides to PPTX converter using Playwright and python-pptx
"""
import asyncio
import os
import sys
from pathlib import Path

# Install playwright if needed
try:
    from playwright.async_api import async_playwright
except ImportError:
    print("Installing playwright...")
    os.system("pip install playwright")
    os.system("playwright install chromium")
    from playwright.async_api import async_playwright

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

WORKING_DIR = Path(__file__).parent
SLIDES_DIR = WORKING_DIR / "slides"
SCREENSHOTS_DIR = WORKING_DIR / "screenshots"
OUTPUT_FILE = WORKING_DIR / "스마트_물류관리_시스템_수행계획서.pptx"

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)  # 1920px at 144dpi
SLIDE_HEIGHT = Inches(7.5)    # 1080px at 144dpi


async def capture_screenshots():
    """Capture screenshots from HTML slides using Playwright"""
    SCREENSHOTS_DIR.mkdir(exist_ok=True)

    html_files = sorted(SLIDES_DIR.glob("slide-*.html"))
    print(f"Found {len(html_files)} HTML slides")

    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(viewport={"width": 1920, "height": 1080})

        for html_file in html_files:
            slide_num = html_file.stem  # e.g., "slide-01-cover"
            output_path = SCREENSHOTS_DIR / f"{slide_num}.png"

            # Load HTML file
            file_url = f"file://{html_file.absolute()}"
            await page.goto(file_url, wait_until="networkidle")

            # Wait for fonts to load
            await page.wait_for_timeout(500)

            # Screenshot the slide element or full page
            try:
                slide = await page.query_selector(".slide")
                if slide:
                    await slide.screenshot(path=str(output_path))
                else:
                    await page.screenshot(path=str(output_path))
            except Exception as e:
                # Fallback to full page screenshot
                await page.screenshot(path=str(output_path), full_page=False)

            print(f"  Captured: {slide_num}")

        await browser.close()

    return sorted(SCREENSHOTS_DIR.glob("slide-*.png"))


def create_pptx(screenshot_files):
    """Create PPTX from screenshot images"""
    prs = Presentation()

    # Set slide dimensions (16:9)
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Blank layout
    blank_layout = prs.slide_layouts[6]  # Blank layout

    for img_path in screenshot_files:
        slide = prs.slides.add_slide(blank_layout)

        # Add image to cover the entire slide
        slide.shapes.add_picture(
            str(img_path),
            Inches(0),
            Inches(0),
            width=SLIDE_WIDTH,
            height=SLIDE_HEIGHT
        )

        print(f"  Added slide: {img_path.name}")

    prs.save(str(OUTPUT_FILE))
    print(f"\nPPTX saved to: {OUTPUT_FILE}")
    return OUTPUT_FILE


async def main():
    print("Stage 5: Generating PPTX")
    print("-" * 40)

    print("\n1. Capturing screenshots from HTML slides...")
    screenshot_files = await capture_screenshots()

    print(f"\n2. Creating PPTX with {len(screenshot_files)} slides...")
    output_path = create_pptx(screenshot_files)

    print("\nDone!")
    return output_path


if __name__ == "__main__":
    asyncio.run(main())
